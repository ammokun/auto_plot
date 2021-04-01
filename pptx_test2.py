from pptx import Presentation
from PIL import Image
import os


SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000
IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH/2, SLIDE_HEIGHT/2
SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT

OUTPUT_FILE_PATH = "test.pptx"
IMG_DIR = "./temp_fig/"

def add_slide(prs):
  blank_slide_layout = prs.slide_layouts[6]    
  slide = prs.slides.add_slide(blank_slide_layout)
  return slide

def add_picture(slide, img_file):
  im = Image.open(img_file)
  im_width, im_height = im.size
  aspect_ratio = im_width/im_height

  if aspect_ratio > SLIDE_ASPECT_RATIO:
    img_display_width = SLIDE_WIDTH
    img_display_height = img_display_width / aspect_ratio
  else:
    img_display_height = SLIDE_HEIGHT
    img_display_width = img_display_height * aspect_ratio
  left = IMG_CENTER_X - img_display_width / 2
  top = IMG_CENTER_Y - img_display_height / 2

  if aspect_ratio > SLIDE_ASPECT_RATIO:
    slide.shapes.add_picture(img_file, left, top, width = img_display_width)  
  else:
    slide.shapes.add_picture(img_file, left, top, height = img_display_height)  

  return slide



prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

img_files = os.listdir(IMG_DIR)
img_files = [name for name in img_files if name.endswith(".eps")]
img_files.sort()
for name in img_files:
  path = IMG_DIR + name
  slide = add_slide(prs)
  add_picture(slide, path)
prs.save(OUTPUT_FILE_PATH)
